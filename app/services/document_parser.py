"""
UniHR 多格式文件解析引擎 (Document Parser Engine)

支援格式：
  Phase 0: PDF(文字型)、DOCX、TXT
  Phase 1: PDF(掃描/OCR)、PDF(表格)、Excel(.xlsx/.xls)、CSV、HTML、Markdown
  Phase 2: RTF、JSON、圖片(JPG/PNG/TIFF)、DOC(舊格式)
  Phase 3: PPT/PPTX、網頁 URL 擷取
  LlamaParse: PDF/DOCX/PPTX/圖片 高品質解析（跨頁表格、手寫 OCR、複雜佈局）

特點：
  - LlamaParse 優先解析（複雜文件品質大幅提升）
  - 精確 Token 計算 (tiktoken)
  - 智慧切片（章節邊界偵測、表格保護）
  - 品質報告系統
  - 多編碼自動偵測
  - OCR 降級策略（pytesseract → LlamaParse Agentic OCR）
  - 網頁 URL 內容擷取（trafilatura）
"""

import os
import re
import csv
import json
import time
import logging
from io import StringIO
from typing import List, Tuple, Dict, Any, Optional
from pathlib import Path
from dataclasses import dataclass, field, asdict

# ── 必要依賴 ──
import pypdf
from docx import Document as DocxDocument
from bs4 import BeautifulSoup

# ── 可選依賴（graceful degradation） ──
try:
    import tiktoken
    _HAS_TIKTOKEN = True
except ImportError:
    _HAS_TIKTOKEN = False

try:
    import pdfplumber
    _HAS_PDFPLUMBER = True
except ImportError:
    _HAS_PDFPLUMBER = False

try:
    import openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False

try:
    import pytesseract
    from PIL import Image
    _HAS_OCR = True
except ImportError:
    _HAS_OCR = False

try:
    from pdf2image import convert_from_path
    _HAS_PDF2IMAGE = True
except ImportError:
    _HAS_PDF2IMAGE = False

try:
    import chardet
    _HAS_CHARDET = True
except ImportError:
    _HAS_CHARDET = False

try:
    from striprtf.striprtf import rtf_to_text
    _HAS_RTF = True
except ImportError:
    _HAS_RTF = False

try:
    from pptx import Presentation as PptxPresentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    import trafilatura
    _HAS_TRAFILATURA = True
except ImportError:
    _HAS_TRAFILATURA = False

# LlamaParse: 延遲導入（lazy import）避免拖慢啟動
# llama_parse 依賴 llama_index_core 很重，只在真正需要時才 import
_HAS_LLAMAPARSE = False
_LlamaParseClient = None

def _get_available_ocr_langs() -> List[str]:
    if not _HAS_OCR:
        return []
    try:
        return pytesseract.get_languages(config="")
    except Exception:
        return []

def _pick_ocr_langs(preferred: str) -> Tuple[str, Optional[str]]:
    langs = _get_available_ocr_langs()
    if not langs:
        return preferred, None

    preferred_list = [p for p in preferred.split("+") if p]
    chosen = [p for p in preferred_list if p in langs]

    if not chosen:
        for cand in ("chi_tra", "chi_sim", "eng"):
            if cand in langs:
                chosen = [cand]
                break

    if not chosen:
        chosen = [sorted(langs)[0]]

    missing = [p for p in preferred_list if p not in chosen]
    if missing:
        note = f"OCR 語言包缺少: {', '.join(missing)}，改用 {'+'.join(chosen)}"
    else:
        note = None

    return "+".join(chosen), note

def _ensure_llamaparse():
    """延遲載入 LlamaParse，只在第一次呼叫時 import"""
    global _HAS_LLAMAPARSE, _LlamaParseClient
    if _LlamaParseClient is not None:
        return True
    try:
        from llama_parse import LlamaParse
        _LlamaParseClient = LlamaParse
        _HAS_LLAMAPARSE = True
        return True
    except ImportError:
        _HAS_LLAMAPARSE = False
        return False

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════
# 品質報告
# ═══════════════════════════════════════════════════════════

@dataclass
class QualityReport:
    """文件解析品質報告"""
    format_detected: str = ""
    total_pages: int = 0
    total_chars: int = 0
    total_chunks: int = 0
    quality_score: float = 0.0      # 0.0 ~ 1.0
    quality_level: str = "unknown"  # excellent / good / fair / poor / failed
    warnings: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)
    suggestions: List[str] = field(default_factory=list)
    tables_detected: int = 0
    images_detected: int = 0
    ocr_used: bool = False
    ocr_confidence: float = 0.0
    encoding_detected: str = ""
    parse_time_ms: int = 0
    parse_engine: str = "native"  # native / llamaparse
    handwriting_detected: bool = False

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)

    def add_warning(self, msg: str):
        self.warnings.append(msg)

    def add_error(self, msg: str):
        self.errors.append(msg)

    def add_suggestion(self, msg: str):
        self.suggestions.append(msg)

    def compute_quality(self):
        """根據各項指標計算品質分數"""
        score = 1.0
        if self.total_chars < 50:
            score -= 0.5
        elif self.total_chars < 200:
            score -= 0.2
        score -= len(self.warnings) * 0.08
        score -= len(self.errors) * 0.3
        if self.ocr_used and self.ocr_confidence < 0.7:
            score -= 0.2
        score = max(0.0, min(1.0, score))
        self.quality_score = round(score, 2)

        if score >= 0.9:
            self.quality_level = "excellent"
        elif score >= 0.7:
            self.quality_level = "good"
        elif score >= 0.5:
            self.quality_level = "fair"
        elif score >= 0.2:
            self.quality_level = "poor"
        else:
            self.quality_level = "failed"


# ═══════════════════════════════════════════════════════════
# 格式映射
# ═══════════════════════════════════════════════════════════

SUPPORTED_FORMATS: Dict[str, str] = {
    # Phase 0
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",
    ".txt": "txt",
    # Phase 1
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".csv": "csv",
    ".html": "html",
    ".htm": "html",
    ".md": "markdown",
    ".markdown": "markdown",
    # Phase 2
    ".rtf": "rtf",
    ".json": "json",
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".tiff": "image",
    ".tif": "image",
    ".bmp": "image",
    ".webp": "image",
    ".heic": "image",
    # Phase 3
    ".pptx": "pptx",
    ".ppt": "ppt",
}

# LlamaParse 支援的格式（這些格式優先走 LlamaParse 以獲得更好的品質）
LLAMAPARSE_FORMATS: set = {
    "pdf", "docx", "doc", "pptx", "ppt",
    "xlsx", "xls",
    "html", "rtf",
    "image",  # 含手寫 OCR
}


# ═══════════════════════════════════════════════════════════
# DocumentParser 主類
# ═══════════════════════════════════════════════════════════

class DocumentParser:
    """
    多格式文件解析引擎。

    使用方式::

        text, metadata = DocumentParser.parse(file_path, file_type)
        # metadata 包含 QualityReport 所有欄位
    """

    @staticmethod
    def get_supported_extensions() -> set:
        """取得所有支援的副檔名"""
        return set(SUPPORTED_FORMATS.keys())

    @staticmethod
    def detect_file_type(filename: str) -> str:
        """根據副檔名偵測文件類型"""
        ext = Path(filename).suffix.lower()
        file_type = SUPPORTED_FORMATS.get(ext)
        if not file_type:
            supported = ", ".join(sorted(SUPPORTED_FORMATS.keys()))
            raise ValueError(
                f"不支援的文件類型: {ext}。支援的類型: {supported}"
            )
        return file_type

    @classmethod
    def parse(cls, file_path: str, file_type: str) -> Tuple[str, dict]:
        """
        解析文件。

        策略：
          1. 若 LlamaParse API Key 可用 且 格式在 LLAMAPARSE_FORMATS 中
             → 優先使用 LlamaParse（品質最高）
          2. LlamaParse 不可用或失敗 → fallback 到內建解析器

        Returns:
            (text_content, metadata_dict)
            metadata_dict 包含 QualityReport 的所有欄位。
        Raises:
            ValueError: 品質為 *failed* 時拋出。
        """
        start = time.time()
        report = QualityReport(format_detected=file_type)

        # ── LlamaParse 優先路徑 ──
        from app.config import settings
        llamaparse_key = getattr(settings, "LLAMAPARSE_API_KEY", "")
        llamaparse_enabled = getattr(settings, "LLAMAPARSE_ENABLED", True)
        if (
            llamaparse_enabled
            and llamaparse_key
            and file_type in LLAMAPARSE_FORMATS
            and _ensure_llamaparse()
        ):
            try:
                text, report = cls._parse_with_llamaparse(
                    file_path, file_type, report, llamaparse_key
                )
                if text.strip() and report.quality_score >= 0.5:
                    report.parse_time_ms = int((time.time() - start) * 1000)
                    report.total_chars = len(text)
                    metadata = report.to_dict()
                    logger.info(
                        f"LlamaParse 解析成功: {file_type}, "
                        f"品質={report.quality_level}, "
                        f"字元={report.total_chars}"
                    )
                    return text.strip(), metadata
                else:
                    logger.warning(
                        f"LlamaParse 解析品質不足 ({report.quality_level})，"
                        f"降級到內建解析器"
                    )
                    # 重置 report 給 fallback 使用
                    report = QualityReport(format_detected=file_type)
            except Exception as e:
                logger.warning(f"LlamaParse 解析失敗: {e}，降級到內建解析器")
                report = QualityReport(format_detected=file_type)

        # ── 內建解析器 ──
        _PARSERS = {
            "pdf": cls._parse_pdf,
            "docx": cls._parse_docx,
            "doc": cls._parse_doc,
            "txt": cls._parse_txt,
            "xlsx": cls._parse_excel,
            "xls": cls._parse_excel,
            "csv": cls._parse_csv,
            "html": cls._parse_html,
            "markdown": cls._parse_markdown,
            "rtf": cls._parse_rtf,
            "json": cls._parse_json,
            "image": cls._parse_image,
            "pptx": cls._parse_pptx,
            "ppt": cls._parse_ppt,
        }

        parser = _PARSERS.get(file_type)
        if not parser:
            raise ValueError(f"不支援的文件類型: {file_type}")

        text, report = parser(file_path, report)

        report.parse_time_ms = int((time.time() - start) * 1000)
        report.total_chars = len(text)
        report.compute_quality()
        metadata = report.to_dict()

        if report.quality_level == "failed":
            detail = "; ".join(report.errors) if report.errors else "品質過低"
            raise ValueError(f"文件解析品質不足: {detail}")

        return text.strip(), metadata

    # ─────────────────────────────────────────────
    # PDF（三層降級：文字 → 表格 → OCR）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_pdf(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        try:
            reader = pypdf.PdfReader(file_path)
            page_count = len(reader.pages)
            report.total_pages = page_count

            # 第一層：文字提取
            text_parts = []
            pages_with_text = 0
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    pages_with_text += 1
                text_parts.append(page_text)
            text = "\n\n".join(text_parts)

            # 掃描偵測（改進版：每頁平均字元數）
            avg_chars = len(text.strip()) / max(page_count, 1)
            scanned_ratio = 1.0 - (pages_with_text / max(page_count, 1))
            is_likely_scanned = avg_chars < 100
            is_partially_scanned = scanned_ratio > 0.3 and not is_likely_scanned

            # 第二層：表格提取（pdfplumber）
            if _HAS_PDFPLUMBER:
                try:
                    table_text, tc = cls._extract_pdf_tables(file_path)
                    if table_text:
                        text += "\n\n" + table_text
                        report.tables_detected = tc
                except Exception as e:
                    report.add_warning(f"表格解析警告: {e}")

            # 第三層：OCR 降級（掃描型 PDF）
            if is_likely_scanned:
                if _HAS_OCR and _HAS_PDF2IMAGE:
                    report.add_warning("偵測到掃描型 PDF，使用 OCR 處理")
                    ocr_text, confidence = cls._ocr_pdf(file_path)
                    if ocr_text.strip():
                        text = ocr_text
                        report.ocr_used = True
                        report.ocr_confidence = confidence
                        if confidence < 0.7:
                            report.add_warning(f"OCR 辨識品質偏低 ({confidence:.0%})，建議上傳文字版 PDF")
                            report.add_suggestion("如果有電子版（Word / 文字型 PDF），請改用電子版以獲得更準確的結果")
                    else:
                        report.add_error("OCR 無法辨識文字內容")
                        report.add_suggestion("請上傳文字型 PDF 或 Word 文件")
                else:
                    report.add_error("偵測到掃描型 PDF，OCR 引擎未安裝")
                    report.add_suggestion("請上傳文字型 PDF 或 Word 文件")
            elif is_partially_scanned:
                report.add_warning(f"部分頁面可能為掃描圖片（{scanned_ratio:.0%} 頁面無文字）")
                report.add_suggestion("部分頁面內容可能遺失，建議檢查結果")

            if not text.strip():
                report.add_error("PDF 文件無法提取任何文字內容")

            return text, report
        except ValueError:
            raise
        except Exception as e:
            report.add_error(f"PDF 解析失敗: {e}")
            return "", report

    @staticmethod
    def _extract_pdf_tables(file_path: str) -> Tuple[str, int]:
        """使用 pdfplumber 提取 PDF 中的表格"""
        table_texts: List[str] = []
        table_count = 0
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                for table in (page.extract_tables() or []):
                    if not table:
                        continue
                    table_count += 1
                    rows = []
                    for row in table:
                        cleaned = [str(c).strip() if c else "" for c in row]
                        rows.append(" | ".join(cleaned))
                    table_texts.append(f"\n[表格 {table_count}]\n" + "\n".join(rows))
        return "\n".join(table_texts), table_count

    @staticmethod
    def _ocr_pdf(file_path: str) -> Tuple[str, float]:
        """OCR 處理掃描型 PDF"""
        from app.config import settings
        images = convert_from_path(file_path, dpi=300)
        all_text: List[str] = []
        confidences: List[float] = []
        lang, _ = _pick_ocr_langs(settings.OCR_LANGS)
        for img in images:
            data = pytesseract.image_to_data(
                img, lang=lang, output_type=pytesseract.Output.DICT
            )
            page_words: List[str] = []
            page_confs: List[float] = []
            for i, word in enumerate(data["text"]):
                conf = int(data["conf"][i])
                if conf > 0 and word.strip():
                    page_words.append(word)
                    page_confs.append(conf / 100.0)
            all_text.append(" ".join(page_words))
            if page_confs:
                confidences.append(sum(page_confs) / len(page_confs))
        avg_conf = sum(confidences) / len(confidences) if confidences else 0.0
        return "\n\n".join(all_text), avg_conf

    # ─────────────────────────────────────────────
    # DOCX（含表格、標題樣式）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_docx(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        try:
            doc = DocxDocument(file_path)
            parts: List[str] = []

            # 段落（保留標題層級）
            for para in doc.paragraphs:
                t = para.text.strip()
                if not t:
                    continue
                style = para.style
                if style and style.name and style.name.startswith("Heading"):
                    level_str = style.name.replace("Heading ", "").replace("Heading", "1")
                    try:
                        level = int(level_str)
                    except ValueError:
                        level = 1
                    parts.append(f"\n{'#' * level} {t}\n")
                else:
                    parts.append(t)

            # 表格
            table_count = 0
            for table in doc.tables:
                table_count += 1
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                parts.append(f"\n[表格 {table_count}]\n" + "\n".join(rows))

            report.tables_detected = table_count
            report.total_pages = max(1, len(doc.paragraphs) // 30)

            text = "\n\n".join(parts)
            if not text.strip():
                report.add_error("DOCX 文件內容為空")
            if table_count > 0:
                report.add_warning(f"偵測到 {table_count} 個表格，已提取內容")

            return text, report
        except Exception as e:
            report.add_error(f"DOCX 解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # DOC 舊格式（antiword / LibreOffice 降級）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_doc(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        import subprocess

        # 嘗試 antiword
        try:
            result = subprocess.run(
                ["antiword", file_path],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                report.add_warning("使用 antiword 解析 .doc 格式，表格格式可能有損失")
                return result.stdout, report
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # 嘗試 libreoffice
        try:
            import tempfile
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "txt",
                     "--outdir", tmpdir, file_path],
                    capture_output=True, text=True, timeout=60,
                )
                if result.returncode == 0:
                    txt_file = os.path.join(tmpdir, Path(file_path).stem + ".txt")
                    if os.path.exists(txt_file):
                        with open(txt_file, "r", encoding="utf-8") as f:
                            text = f.read()
                        report.add_warning("使用 LibreOffice 轉換 .doc 格式")
                        return text, report
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        report.add_error(".doc 格式需要 antiword 或 LibreOffice 來解析")
        report.add_suggestion("請將 .doc 文件另存為 .docx 或 .pdf 格式後重新上傳")
        return "", report

    # ─────────────────────────────────────────────
    # TXT（智慧編碼偵測）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_txt(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        try:
            with open(file_path, "rb") as f:
                raw = f.read()

            # 二進位檔案偵測：NULL bytes 幾乎不存在於文字檔案
            null_count = raw.count(b'\x00')
            if null_count >= 2:
                report.add_error("偵測到二進位檔案內容（包含大量 NULL bytes）")
                report.add_suggestion("請確認上傳的是文字檔案而非二進位檔案")
                return "", report

            # 編碼偵測（大檔案只取前 64KB 偵測以提升效能）
            encoding = "utf-8"
            detect_sample = raw[:65536] if len(raw) > 65536 else raw
            if _HAS_CHARDET:
                detected = chardet.detect(detect_sample)
                if detected and detected.get("encoding"):
                    encoding = detected["encoding"]
                    conf = detected.get("confidence", 0)
                    report.encoding_detected = f"{encoding} ({conf:.0%})"
                    if conf < 0.7:
                        report.add_warning(f"編碼偵測信心度較低 ({conf:.0%})，可能有亂碼")

            # 依序嘗試解碼（utf-8-sig 優先以處理 BOM）
            text = None
            for enc in ["utf-8-sig", encoding, "utf-8", "big5", "gbk", "cp1252", "latin-1"]:
                try:
                    text = raw.decode(enc)
                    if text.strip():
                        break
                except (UnicodeDecodeError, LookupError):
                    continue

            if text is None or not text.strip():
                report.add_error("無法解碼文字檔案內容")
                return "", report

            # 清除 BOM 字元
            if text.startswith("\ufeff"):
                text = text.lstrip("\ufeff")

            # 亂碼 / 二進位偵測
            sample = text[:2000]
            garbled = sum(1 for c in sample if ord(c) > 0xFFF0) / max(len(sample), 1)
            # 控制字元偵測（排除 \t\n\r）
            control_chars = sum(1 for c in sample if ord(c) < 32 and c not in '\t\n\r') / max(len(sample), 1)
            if garbled > 0.1 or control_chars > 0.1:
                report.add_warning("偵測到可能的亂碼或二進位內容，建議以 UTF-8 編碼儲存後重新上傳")
            # 可讀性偵測：判斷是否為真實文字內容
            # 真實文字應包含字母 / CJK / 數字 / 空白，而非隨機二進位解碼出的字元
            text_like_chars = sum(1 for c in sample if (
                c.isalpha() or c.isdigit() or c.isspace() or
                "\u4e00" <= c <= "\u9fff" or  # CJK
                c in '.,;:!?\'"-()[]{}\n\t/\\@#$%&*+=<>~`'
            )) / max(len(sample), 1)
            if text_like_chars < 0.4:
                report.add_error("檔案內容大部分為不可讀字元（可能為二進位檔案）")
                report.add_suggestion("請確認上傳的是文字檔案而非二進位檔案")

            report.total_pages = max(1, len(text.split("\n")) // 50)
            return text, report
        except Exception as e:
            report.add_error(f"TXT 解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # Excel (.xlsx / .xls)
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_excel(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        if not _HAS_OPENPYXL:
            report.add_error("Excel 解析引擎未安裝 (openpyxl)")
            report.add_suggestion("請將 Excel 文件轉為 PDF 或 CSV 後上傳")
            return "", report

        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts: List[str] = []
            table_count = 0

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    report.add_warning(f"工作表「{sheet_name}」為空")
                    continue

                valid_rows = []
                for row in rows:
                    cleaned = [str(c).strip() if c is not None else "" for c in row]
                    if any(cleaned):
                        valid_rows.append(cleaned)

                if not valid_rows:
                    report.add_warning(f"工作表「{sheet_name}」無有效資料")
                    continue

                table_count += 1
                sheet_text = f"\n## 工作表：{sheet_name}\n"
                for i, row in enumerate(valid_rows):
                    if i == 0:
                        sheet_text += " | ".join(row) + "\n"
                        sheet_text += " | ".join(["---"] * len(row)) + "\n"
                    else:
                        sheet_text += " | ".join(row) + "\n"
                parts.append(sheet_text)

            wb.close()
            report.tables_detected = table_count
            report.total_pages = table_count

            if not parts:
                report.add_error("Excel 文件無可讀取的內容")
                return "", report

            report.add_warning(f"偵測到 {table_count} 個工作表資料")
            return "\n".join(parts), report
        except Exception as e:
            report.add_error(f"Excel 解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # CSV
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_csv(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        try:
            with open(file_path, "rb") as f:
                raw = f.read()

            encoding = "utf-8"
            if _HAS_CHARDET:
                detected = chardet.detect(raw)
                if detected and detected.get("encoding"):
                    encoding = detected["encoding"]

            text_content = raw.decode(encoding, errors="replace")

            try:
                dialect = csv.Sniffer().sniff(text_content[:4096])
                delimiter = dialect.delimiter
            except csv.Error:
                delimiter = ","

            reader = csv.reader(StringIO(text_content), delimiter=delimiter)
            rows = list(reader)
            if not rows:
                report.add_error("CSV 文件為空")
                return "", report

            parts = []
            for i, row in enumerate(rows):
                if i == 0:
                    parts.append(" | ".join(row))
                    parts.append(" | ".join(["---"] * len(row)))
                else:
                    parts.append(" | ".join(row))

            report.tables_detected = 1
            report.total_pages = 1
            return "\n".join(parts), report
        except Exception as e:
            report.add_error(f"CSV 解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # HTML
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_html(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        try:
            with open(file_path, "rb") as f:
                raw = f.read()

            encoding = "utf-8"
            if _HAS_CHARDET:
                detected = chardet.detect(raw)
                if detected and detected.get("encoding"):
                    encoding = detected["encoding"]

            html = raw.decode(encoding, errors="replace")
            soup = BeautifulSoup(html, "lxml")

            for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                tag.decompose()

            parts: List[str] = []
            for el in soup.find_all(
                ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th", "pre", "blockquote"]
            ):
                tag_name = el.name
                t = el.get_text(strip=True)
                if not t:
                    continue
                if tag_name.startswith("h"):
                    level = int(tag_name[1])
                    parts.append(f"\n{'#' * level} {t}\n")
                elif tag_name == "li":
                    parts.append(f"- {t}")
                elif tag_name in ("pre", "blockquote"):
                    parts.append(f"\n> {t}\n")
                else:
                    parts.append(t)

            report.tables_detected = len(soup.find_all("table"))
            result = "\n\n".join(parts) if parts else soup.get_text(separator="\n", strip=True)

            if not result.strip():
                report.add_error("HTML 文件無可讀取的文字內容")

            report.total_pages = 1
            return result, report
        except Exception as e:
            report.add_error(f"HTML 解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # Markdown
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_markdown(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            if not text.strip():
                report.add_error("Markdown 文件為空")
            report.total_pages = max(1, len(text.split("\n")) // 50)
            return text, report
        except UnicodeDecodeError:
            return cls._parse_txt(file_path, report)
        except Exception as e:
            report.add_error(f"Markdown 解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # RTF
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_rtf(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        if not _HAS_RTF:
            report.add_error("RTF 解析引擎未安裝 (striprtf)")
            report.add_suggestion("請將 RTF 文件轉為 DOCX 或 PDF 後上傳")
            return "", report
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                rtf_content = f.read()
            text = rtf_to_text(rtf_content)
            if not text.strip():
                report.add_error("RTF 文件內容為空")
            report.total_pages = max(1, len(text.split("\n")) // 50)
            return text, report
        except Exception as e:
            report.add_error(f"RTF 解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # JSON（結構化資料 → 可讀文字）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_json(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            text = cls._json_to_text(data)
            if not text.strip():
                report.add_error("JSON 文件無可讀取的內容")
            report.total_pages = 1
            return text, report
        except json.JSONDecodeError as e:
            report.add_error(f"JSON 格式錯誤: {e}")
            report.compute_quality()
            raise ValueError(f"JSON 格式錯誤: {e}")
        except ValueError:
            raise
        except Exception as e:
            report.add_error(f"JSON 解析失敗: {e}")
            return "", report

    @staticmethod
    def _json_to_text(data: Any, prefix: str = "") -> str:
        lines: List[str] = []
        if isinstance(data, dict):
            for k, v in data.items():
                if isinstance(v, (dict, list)):
                    lines.append(f"{prefix}{k}:")
                    lines.append(DocumentParser._json_to_text(v, prefix + "  "))
                else:
                    lines.append(f"{prefix}{k}: {v}")
        elif isinstance(data, list):
            for i, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    lines.append(f"{prefix}[{i + 1}]")
                    lines.append(DocumentParser._json_to_text(item, prefix + "  "))
                else:
                    lines.append(f"{prefix}- {item}")
        else:
            lines.append(f"{prefix}{data}")
        return "\n".join(lines)

    # ─────────────────────────────────────────────
    # 圖片（OCR）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_image(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        if not _HAS_OCR:
            report.add_error("圖片 OCR 引擎未安裝 (pytesseract)")
            report.add_suggestion("請將圖片中的文字轉為文字檔後上傳")
            return "", report
        try:
            img = Image.open(file_path)
            report.images_detected = 1
            report.total_pages = 1

            from app.config import settings
            lang, note = _pick_ocr_langs(settings.OCR_LANGS)
            if note:
                report.add_warning(note)

            data = pytesseract.image_to_data(
                img, lang=lang, output_type=pytesseract.Output.DICT
            )
            words: List[str] = []
            confs: List[float] = []
            for i, word in enumerate(data["text"]):
                conf = int(data["conf"][i])
                if conf > 0 and word.strip():
                    words.append(word)
                    confs.append(conf / 100.0)

            text = " ".join(words)
            report.ocr_used = True
            report.ocr_confidence = sum(confs) / len(confs) if confs else 0.0

            if not text.strip():
                report.add_error("圖片 OCR 無法辨識文字")
            elif report.ocr_confidence < 0.7:
                report.add_warning(f"OCR 辨識信心度偏低 ({report.ocr_confidence:.0%})")

            return text, report
        except Exception as e:
            report.add_error(f"圖片解析失敗: {e}")
            return "", report

    # ─────────────────────────────────────────────
    # PPT/PPTX（投影片 + 表格 + 備忘稿）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_pptx(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        if not _HAS_PPTX:
            report.add_error("PPTX 解析引擎未安裝 (python-pptx)")
            report.add_suggestion("請將 PPTX 轉為 PDF 後上傳")
            return "", report
        try:
            prs = PptxPresentation(file_path)
            parts: List[str] = []
            table_count = 0

            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_parts: List[str] = [f"\n## 投影片 {slide_idx}\n"]

                for shape in slide.shapes:
                    # 文字框
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                # 偵測標題（通常是較大字體或 Title 佔位符）
                                if shape.shape_type is not None and hasattr(shape, "placeholder_format"):
                                    if shape.placeholder_format and shape.placeholder_format.idx in (0, 1):
                                        slide_parts.append(f"### {t}")
                                        continue
                                slide_parts.append(t)

                    # 表格
                    if shape.has_table:
                        table_count += 1
                        rows = []
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append(" | ".join(cells))
                        if rows:
                            # 自動加表頭分隔線
                            table_text = rows[0] + "\n"
                            table_text += " | ".join(["---"] * len(shape.table.rows[0].cells)) + "\n"
                            table_text += "\n".join(rows[1:])
                            slide_parts.append(f"\n[表格 {table_count}]\n{table_text}")

                # 備忘稿（speaker notes）
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"\n> 備忘稿：{notes}")

                parts.append("\n".join(slide_parts))

            report.total_pages = len(prs.slides)
            report.tables_detected = table_count
            text = "\n\n".join(parts)

            if not text.strip():
                report.add_error("PPTX 文件無可讀取的內容")
            if table_count > 0:
                report.add_warning(f"偵測到 {table_count} 個投影片表格")

            return text, report
        except Exception as e:
            report.add_error(f"PPTX 解析失敗: {e}")
            return "", report

    @classmethod
    def _parse_ppt(cls, file_path: str, report: QualityReport) -> Tuple[str, QualityReport]:
        """舊格式 .ppt — 嘗試 LibreOffice 轉換，或提示用戶轉檔"""
        import subprocess
        try:
            import tempfile
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "pptx",
                     "--outdir", tmpdir, file_path],
                    capture_output=True, text=True, timeout=120,
                )
                if result.returncode == 0:
                    pptx_file = os.path.join(tmpdir, Path(file_path).stem + ".pptx")
                    if os.path.exists(pptx_file):
                        report.add_warning("使用 LibreOffice 將 .ppt 轉換為 .pptx 進行解析")
                        return cls._parse_pptx(pptx_file, report)
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        report.add_error(".ppt 格式需要 LibreOffice 來解析")
        report.add_suggestion("請將 .ppt 文件另存為 .pptx 或 .pdf 格式後重新上傳")
        return "", report

    # ─────────────────────────────────────────────
    # LlamaParse 高品質解析（跨頁表格、手寫 OCR、複雜佈局）
    # ─────────────────────────────────────────────
    @classmethod
    def _parse_with_llamaparse(
        cls,
        file_path: str,
        file_type: str,
        report: QualityReport,
        api_key: str,
    ) -> Tuple[str, QualityReport]:
        """
        使用 LlamaParse 進行高品質解析。

        優點：
          - 跨頁表格完整保留
          - 手寫文字 OCR（Agentic OCR）
          - 複雜佈局（雙欄、嵌套表格）
          - 圖片中的文字提取

        若解析品質不足或 API 不可用，呼叫者應 fallback 到內建解析器。
        """
        import nest_asyncio
        nest_asyncio.apply()

        # 根據文件類型調整解析參數
        parsing_instruction = (
            "這是一份人力資源或企業管理相關的文件。"
            "請精確保留所有表格結構（使用 Markdown 表格格式）、"
            "包含表頭的對應關係。"
            "如果有手寫文字，請盡力辨識。"
            "保留文件的標題層級結構。"
            "輸出語言與原文件一致。"
        )

        # 是否為圖片（啟用更強的 OCR）
        is_image = file_type == "image"

        from app.config import settings

        def run_parse(params: Dict[str, Any]) -> List[Any]:
            parser = _LlamaParseClient(**params)
            import asyncio
            try:
                loop = asyncio.get_event_loop()
            except RuntimeError:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
            return loop.run_until_complete(parser.aload_data(file_path))

        try:
            params = {
                "api_key": api_key,
                "result_type": settings.LLAMAPARSE_RESULT_TYPE,
                "parsing_instruction": parsing_instruction,
                "language": settings.LLAMAPARSE_LANGUAGE,
            }
            if settings.LLAMAPARSE_AUTO_MODE:
                params.update({
                    "auto_mode": True,
                    "auto_mode_trigger_on_image_in_page": True,
                    "auto_mode_trigger_on_table_in_page": True,
                })

            try:
                documents = run_parse(params)
            except Exception as e:
                msg = str(e)
                if "422" in msg or "Unprocessable" in msg:
                    report.add_warning("LlamaParse 參數疑似不相容，改用精簡參數重試")
                    minimal = {
                        "api_key": api_key,
                        "result_type": settings.LLAMAPARSE_RESULT_TYPE,
                        "parsing_instruction": parsing_instruction,
                    }
                    documents = run_parse(minimal)
                else:
                    raise

            if not documents:
                report.add_warning("LlamaParse 返回空結果")
                report.compute_quality()
                return "", report

            # 合併所有頁面
            text_parts = []
            for doc in documents:
                if hasattr(doc, "text") and doc.text:
                    text_parts.append(doc.text)

            text = "\n\n".join(text_parts)

            # 統計
            report.parse_engine = "llamaparse"
            report.total_pages = len(documents)
            report.total_chars = len(text)

            # 偵測表格
            table_matches = re.findall(r"\|.*\|", text)
            if table_matches:
                # 粗估表格數量（連續的表格行視為一個表格）
                table_lines = [i for i, line in enumerate(text.split("\n")) if "|" in line and line.strip().startswith("|")]
                if table_lines:
                    table_groups = 1
                    for i in range(1, len(table_lines)):
                        if table_lines[i] - table_lines[i-1] > 2:
                            table_groups += 1
                    report.tables_detected = table_groups

            # 偵測是否使用了 OCR（圖片檔或掃描類）
            if is_image:
                report.ocr_used = True
                report.ocr_confidence = 0.85  # LlamaParse OCR 品質通常較高
                report.handwriting_detected = True

            report.compute_quality()
            return text, report

        except Exception as e:
            report.add_warning(f"LlamaParse 錯誤: {e}")
            report.compute_quality()
            return "", report

    # ─────────────────────────────────────────────
    # 網頁 URL 擷取
    # ─────────────────────────────────────────────
    @staticmethod
    def parse_url(url: str) -> Tuple[str, dict]:
        """
        從 URL 擷取網頁正文內容。

        使用 trafilatura 進行高品質網頁正文提取，
        自動移除導覽列、頁尾、廣告等非正文內容。

        Args:
            url: 網頁 URL

        Returns:
            (text_content, metadata_dict)
        """
        start = time.time()
        report = QualityReport(format_detected="url")

        if not _HAS_TRAFILATURA:
            report.add_error("網頁擷取引擎未安裝 (trafilatura)")
            report.compute_quality()
            raise ValueError("trafilatura 未安裝，無法擷取網頁")

        try:
            # 下載網頁
            downloaded = trafilatura.fetch_url(url)
            if not downloaded:
                report.add_error(f"無法下載網頁: {url}")
                report.compute_quality()
                raise ValueError(f"無法下載網頁: {url}")

            # 提取正文（含表格、標題層級）
            text = trafilatura.extract(
                downloaded,
                include_tables=True,
                include_links=False,
                include_images=False,
                include_comments=False,
                output_format="txt",
                favor_precision=True,
            )

            if not text or not text.strip():
                report.add_error("網頁內容為空或無法提取正文")
                report.compute_quality()
                raise ValueError("網頁內容為空")

            # 嘗試同時取得 metadata
            metadata_extracted = trafilatura.extract(
                downloaded,
                output_format="xml",
                include_tables=True,
            )

            report.total_chars = len(text)
            report.total_pages = 1
            report.parse_engine = "trafilatura"
            report.parse_time_ms = int((time.time() - start) * 1000)
            report.compute_quality()

            return text.strip(), report.to_dict()

        except ValueError:
            raise
        except Exception as e:
            report.add_error(f"網頁擷取失敗: {e}")
            report.compute_quality()
            raise ValueError(f"網頁擷取失敗: {e}")


# ═══════════════════════════════════════════════════════════
# TextChunker — 智慧切片
# ═══════════════════════════════════════════════════════════

class TextChunker:
    """
    智慧文字切片器

    - 精確 Token 計算（tiktoken，無可用時使用估算）
    - 章節邊界偵測（Markdown 標題、分隔線）
    - 表格保護（「[表格 N]」區塊不拆開）
    - 句子級強制切分（超長段落）
    - 重疊區保留上下文
    """

    _encoder = None

    @classmethod
    def _get_encoder(cls):
        if cls._encoder is None and _HAS_TIKTOKEN:
            cls._encoder = tiktoken.encoding_for_model("gpt-3.5-turbo")
        return cls._encoder

    @classmethod
    def count_tokens(cls, text: str) -> int:
        """精確計算 token 數量"""
        encoder = cls._get_encoder()
        if encoder:
            return len(encoder.encode(text))
        # 估算：中文 1 字 ≈ 1.5 tokens
        cn = sum(1 for c in text if "\u4e00" <= c <= "\u9fff")
        en_words = len(re.findall(r"[a-zA-Z]+", text))
        other = len(text) - cn - sum(len(w) for w in re.findall(r"[a-zA-Z]+", text))
        return int(cn * 1.5 + en_words + other * 0.5)

    @classmethod
    def split_by_tokens(
        cls,
        text: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 150,
    ) -> List[str]:
        """
        將文本切片為固定 Token 大小的片段。

        1. 按章節 / 段落邊界切分
        2. 保護表格不被拆開
        3. 精確 Token 計算
        4. 重疊保留上下文
        """
        if not text.strip():
            return []

        sections = cls._split_into_sections(text)
        # 確保空字串 section 不會被累積
        sections = [s for s in sections if s.strip()]
        from app.config import settings
        min_section_tokens = (
            settings.MARKDOWN_MIN_SECTION_TOKENS
            if cls._is_markdown_like(text)
            else settings.TEXT_MIN_SECTION_TOKENS
        )
        sections = cls._merge_small_sections(sections, min_section_tokens)
        chunks: List[str] = []
        current_chunk = ""
        current_tokens = 0

        for section in sections:
            section = section.strip()
            if not section:
                continue

            section_tokens = cls.count_tokens(section)

            # 章節標題開頭 → 若當前 chunk 已有足夠內容則強制斷開
            is_heading = section.lstrip().startswith("#")
            if is_heading and current_chunk.strip() and current_tokens >= 30:
                chunks.append(current_chunk.strip())
                current_chunk = section
                current_tokens = section_tokens
                continue

            # 超大段落 → 強制句子級拆分
            if section_tokens > chunk_size:
                if current_chunk.strip():
                    chunks.append(current_chunk.strip())
                sub_chunks = cls._force_split(section, chunk_size, chunk_overlap)
                chunks.extend(sub_chunks)
                current_chunk = ""
                current_tokens = 0
                continue

            # 累積直到超過 chunk_size
            if current_tokens + section_tokens > chunk_size:
                if current_chunk.strip():
                    chunks.append(current_chunk.strip())
                overlap = cls._get_overlap(current_chunk, chunk_overlap)
                current_chunk = (overlap + "\n\n" + section) if overlap else section
                current_tokens = cls.count_tokens(current_chunk)
            else:
                current_chunk += ("\n\n" + section) if current_chunk else section
                current_tokens += section_tokens

        if current_chunk.strip():
            chunks.append(current_chunk.strip())

        # 過濾過小碎片
        return [c for c in chunks if cls.count_tokens(c) >= settings.TEXT_MIN_SECTION_TOKENS]

    # ─── 輔助 ───

    @staticmethod
    def _split_into_sections(text: str) -> List[str]:
        """按章節邊界切分，保護表格區塊不被拆開"""
        # 保護表格
        table_re = re.compile(r"\[表格 \d+\][\s\S]*?(?=\n\n|\[表格|\Z)")
        tables: Dict[str, str] = {}
        for i, m in enumerate(table_re.finditer(text)):
            ph = f"__TBL_{i}__"
            tables[ph] = m.group()
        protected = text
        for ph, tbl in tables.items():
            protected = protected.replace(tbl, ph)

        # 優先 Markdown 標題切分（包含文件開頭的標題）
        heading_re = re.compile(r"(?:^|\n)(?=#{1,6}\s)")
        if heading_re.search(protected):
            sections = heading_re.split(protected)
        else:
            sections = protected.split("\n\n")

        # 還原表格
        result: List[str] = []
        for sec in sections:
            for ph, tbl in tables.items():
                sec = sec.replace(ph, tbl)
            result.append(sec)
        return result

    @classmethod
    def _merge_small_sections(cls, sections: List[str], min_tokens: int) -> List[str]:
        if min_tokens <= 0:
            return sections
        merged: List[str] = []
        buffer = ""
        buffer_tokens = 0
        for sec in sections:
            sec = sec.strip()
            if not sec:
                continue
            sec_tokens = cls.count_tokens(sec)
            if buffer:
                if buffer_tokens + sec_tokens < min_tokens:
                    buffer = buffer + "\n\n" + sec
                    buffer_tokens = cls.count_tokens(buffer)
                    continue
                merged.append(buffer)
                buffer = ""
                buffer_tokens = 0

            if sec_tokens < min_tokens:
                buffer = sec
                buffer_tokens = sec_tokens
            else:
                merged.append(sec)

        if buffer:
            merged.append(buffer)
        return merged

    @staticmethod
    def _is_markdown_like(text: str) -> bool:
        lines = text.splitlines()[:200]
        if not lines:
            return False
        heading_lines = sum(1 for l in lines if l.lstrip().startswith("#"))
        return heading_lines >= 3 or heading_lines / max(len(lines), 1) > 0.1

    @classmethod
    def _force_split(cls, text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """超長段落的句子級強制拆分"""
        sentences = re.split(r"(?<=[。！？\.\!\?\n])", text)
        chunks: List[str] = []
        current = ""
        current_tokens = 0

        for sent in sentences:
            st = cls.count_tokens(sent)
            if current_tokens + st > chunk_size:
                if current.strip():
                    chunks.append(current.strip())
                overlap = cls._get_overlap(current, chunk_overlap)
                current = (overlap + sent) if overlap else sent
                current_tokens = cls.count_tokens(current)
            else:
                current += sent
                current_tokens += st

        if current.strip():
            chunks.append(current.strip())
        return chunks

    @classmethod
    def _get_overlap(cls, text: str, overlap_tokens: int) -> str:
        """取 chunk 尾部作為重疊區域"""
        if not text:
            return ""
        encoder = cls._get_encoder()
        if encoder:
            tokens = encoder.encode(text)
            if len(tokens) <= overlap_tokens:
                return text
            return encoder.decode(tokens[-overlap_tokens:])
        # 估算
        char_est = overlap_tokens * 2
        return text[-char_est:] if len(text) > char_est else text
